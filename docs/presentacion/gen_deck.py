# -*- coding: utf-8 -*-
"""
Genera DECK_GENESIS.pptx (16:9) a partir del guion de slides de GENESIS.
Tema: JARVIS oscuro con acentos aguamarina (#2dffae). Una lamina por "## Slide N".
Contenido FIEL al markdown fuente.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "DECK_GENESIS.pptx")

# ---- Paleta ----
BG       = RGBColor(0x0A, 0x1A, 0x25)   # fondo oscuro petroleo
BG2      = RGBColor(0x0D, 0x22, 0x30)   # paneles
ACCENT   = RGBColor(0x2D, 0xFF, 0xAE)   # aguamarina
ACCENT_D = RGBColor(0x18, 0xC9, 0x8A)
WHITE    = RGBColor(0xF4, 0xFB, 0xF8)
MUTE     = RGBColor(0x9F, 0xC7, 0xBB)   # gris-azulado
GRID     = RGBColor(0x12, 0x31, 0x3F)
TBL_HD   = RGBColor(0x07, 0x16, 0x1F)
TBL_R1   = RGBColor(0x0E, 0x26, 0x33)
TBL_R2   = RGBColor(0x10, 0x2E, 0x3C)

EMUW = Inches(13.333)
EMUH = Inches(7.5)
FONT = "Segoe UI"
FONT_B = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width = EMUW
prs.slide_height = EMUH
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    # fondo
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMUW, EMUH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    _send_back(s, bg)
    return s


def _send_back(slide, shape):
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _no_shadow(shape):
    shape.shadow.inherit = False


def accent_bar(s, x=Inches(0.6), y=Inches(1.55), w=Inches(1.4), h=Pt(4)):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); _no_shadow(bar)
    return bar


def footer(s, idx, total):
    # linea inferior
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.02),
                            Inches(12.13), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = GRID
    ln.line.fill.background(); _no_shadow(ln)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(6), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "GENESIS  ·  Asistente de IA personal, 100% local"
    r.font.size = Pt(9); r.font.color.rgb = MUTE; r.font.name = FONT
    tb2 = s.shapes.add_textbox(Inches(11.2), Inches(7.05), Inches(1.53), Inches(0.35))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = "%02d / %02d" % (idx, total)
    r2.font.size = Pt(9); r2.font.color.rgb = ACCENT_D; r2.font.name = FONT_B


def tech_corners(s):
    L = Inches(0.45); sz = Inches(0.5); th = Pt(2)
    coords = [
        (Inches(0.35), Inches(0.35), 1, 1),
        (Inches(12.5), Inches(0.35), -1, 1),
        (Inches(0.35), Inches(6.65), 1, -1),
        (Inches(12.5), Inches(6.65), -1, -1),
    ]
    for x, y, fx, fy in coords:
        h = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x if fx > 0 else x - sz, y, sz, th)
        h.fill.solid(); h.fill.fore_color.rgb = ACCENT; h.line.fill.background(); _no_shadow(h)
        v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x if fx > 0 else x, y, th, sz if fy > 0 else sz)
        if fy < 0:
            v.top = y - sz + th
        v.fill.solid(); v.fill.fore_color.rgb = ACCENT; v.line.fill.background(); _no_shadow(v)


def title(s, text, sub=None):
    accent_bar(s)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(12.1), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_B
    if sub:
        tb2 = s.shapes.add_textbox(Inches(0.62), Inches(1.62), Inches(12.0), Inches(0.5))
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(15); r2.font.color.rgb = ACCENT; r2.font.name = FONT


def bullets(s, items, top=Inches(2.25), left=Inches(0.75),
            width=Inches(11.8), size=18, gap=10, leader_color=None):
    tb = s.shapes.add_textbox(left, top, width, Inches(4.4))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        # vinieta
        rb = p.add_run(); rb.text = "▸  "
        rb.font.color.rgb = ACCENT; rb.font.size = Pt(size); rb.font.name = FONT_B
        for seg, bold in it if isinstance(it, list) else [(it, False)]:
            r = p.add_run(); r.text = seg
            r.font.size = Pt(size); r.font.name = FONT
            r.font.bold = bold
            r.font.color.rgb = WHITE if bold else MUTE
    return tb


def panel(s, x, y, w, h, color=BG2, line=ACCENT, line_w=Pt(1)):
    pn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pn.fill.solid(); pn.fill.fore_color.rgb = color
    pn.line.color.rgb = line; pn.line.width = line_w
    _no_shadow(pn)
    try:
        pn.adjustments[0] = 0.04
    except Exception:
        pass
    return pn


# =================================================================
# SLIDE 1 — Portada
# =================================================================
def slide1(total):
    s = add_slide()
    tech_corners(s)
    # anillos JARVIS (elipses concentricas, sin relleno)
    cx, cy = Inches(6.6665), Inches(3.55)
    for d in (Inches(4.6), Inches(5.4)):
        ring = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - d/2, cy - d/2, d, d)
        ring.fill.background()
        ring.line.color.rgb = RGBColor(0x1F, 0x6E, 0x57); ring.line.width = Pt(1)
        _no_shadow(ring)
    ring2 = s.shapes.add_shape(MSO_SHAPE.ARC, cx - Inches(2.0), cy - Inches(2.0),
                               Inches(4.0), Inches(4.0))
    ring2.fill.background(); ring2.line.color.rgb = ACCENT; ring2.line.width = Pt(2.5)
    _no_shadow(ring2)
    # titulo
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.55), Inches(12.33), Inches(1.6))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "GENESIS"
    r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_B
    # subtitulo
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(4.05), Inches(12.33), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Asistente de IA personal, 100% local"
    r2.font.size = Pt(22); r2.font.color.rgb = ACCENT; r2.font.name = FONT
    # tagline
    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(4.75), Inches(12.33), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "Privado  ·  Sin nube  ·  Soberano"
    r3.font.size = Pt(15); r3.font.italic = True; r3.font.color.rgb = MUTE; r3.font.name = FONT
    # autor
    tb4 = s.shapes.add_textbox(Inches(0.5), Inches(6.25), Inches(12.33), Inches(0.5))
    p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run(); r4.text = "Alex Quiñones — Junio 2026"
    r4.font.size = Pt(13); r4.font.color.rgb = MUTE; r4.font.name = FONT


# =================================================================
# SLIDE 2 — El problema
# =================================================================
def slide2(total):
    s = add_slide()
    title(s, "El problema",
          "Los asistentes de IA hoy te cuestan tu privacidad.")
    bullets(s, [
        [("☁️ Dependen de la nube", True), (" → sin internet no funcionan", False)],
        [("\U0001F513 Tus conversaciones y archivos viajan a terceros", False)],
        [("\U0001F4B3 Suscripciones y cobros por uso", False)],
        [("\U0001F512 Caja negra: no controlas el modelo ni tus datos", False)],
        [("\U0001F5E3️ Solo conversan; no actuan sobre tu equipo", False)],
    ], top=Inches(2.45), size=19, gap=14)
    footer(s, 2, total)


# =================================================================
# SLIDE 3 — La solucion
# =================================================================
def slide3(total):
    s = add_slide()
    title(s, "La solucion: GENESIS",
          "Un asistente que corre entero en tu PC.")
    bullets(s, [
        [("✅ ", False), ("100% local", True), (" — funciona sin internet", False)],
        [("✅ ", False), ("Privado por diseno", True), (" — tus datos nunca salen", False)],
        [("✅ ", False), ("Sin suscripciones", True), (" — no necesita APIs de pago", False)],
        [("✅ ", False), ("Controlas el modelo", True)],
        [("✅ ", False), ("Actua", True), (" sobre tu sistema, no solo conversa", False)],
    ], top=Inches(2.45), size=19, gap=14)
    footer(s, 3, total)


# =================================================================
# SLIDE 4 — Sentidos y creacion
# =================================================================
def slide4(total):
    s = add_slide()
    title(s, "Que hace (sentidos y creacion)",
          "Piensa  ·  Ve  ·  Habla  ·  Crea")
    bullets(s, [
        [("\U0001F9E0 Lenguaje: ", True), ("modelos locales (Ollama) con router y memoria", False)],
        [("\U0001F399️ Voz: ", True), ("22 voces + ", False), ("clonacion de voz", True), (" local", False)],
        [("\U0001F3A8 Imagenes: ", True), ("Stable Diffusion en tu GPU (offline)", False)],
        [("\U0001F441️ Vision: ", True), ("analiza imagenes y pantalla", False)],
    ], top=Inches(2.55), size=20, gap=18)
    footer(s, 4, total)


# =================================================================
# SLIDE 5 — Accion real
# =================================================================
def slide5(total):
    s = add_slide()
    title(s, "Que hace (accion real)",
          "Controla tu computadora de verdad")
    bullets(s, [
        [("\U0001F39B️ Sistema: ", True), ("volumen, brillo, energia, impresion, multi-monitor", False)],
        [("\U0001F50C Conexiones: ", True), ("WiFi, Bluetooth, USB", False)],
        [("\U0001F4C2 Archivos y apps: ", True), ("abrir, buscar, gestionar (papelera segura)", False)],
        [("\U0001F4FA Musica y casting a TV ", True), ("(Chromecast) + streaming", False)],
        [("\U0001F4E7 Correo, despertador con voz, recordatorios", False)],
    ], top=Inches(2.45), size=19, gap=14)
    footer(s, 5, total)


# =================================================================
# SLIDE 6 — Interfaz (con placeholder de captura)
# =================================================================
def slide6(total):
    s = add_slide()
    title(s, "Interfaz", "Una cabina tipo “JARVIS”")
    bullets(s, [
        [("Núcleo de plasma reactivo a la voz", False)],
        [("Dock de acciones rapidas (buscar, web, musica, crear, ver, nucleo)", False)],
        [("Tablero de evidencias + telemetria del sistema en vivo", False)],
        [("App de escritorio + interfaz web local", False)],
    ], top=Inches(2.35), left=Inches(0.75), width=Inches(6.0), size=16, gap=12)
    # placeholder de captura
    ph = panel(s, Inches(7.1), Inches(2.3), Inches(5.6), Inches(4.0),
               color=BG2, line=ACCENT)
    # marca de esquina diagonal
    tf = ph.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "\U0001F5BC️"
    r.font.size = Pt(34); r.font.color.rgb = ACCENT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "(Insertar captura de pantalla del hub aqui.)"
    r2.font.size = Pt(13); r2.font.italic = True; r2.font.color.rgb = MUTE; r2.font.name = FONT
    footer(s, 6, total)


# =================================================================
# SLIDE 7 — El diferencial (tabla)
# =================================================================
def slide7(total):
    s = add_slide()
    title(s, "El diferencial")
    data = [
        ["", "GENESIS", "Comerciales"],
        ["Procesamiento", "Local", "Nube"],
        ["Privacidad", "Total", "Terceros"],
        ["Offline", "Si", "No"],
        ["Costo", "$0 recurrente", "Suscripcion"],
        ["Accion en el equipo", "Amplia", "Limitada"],
    ]
    rows, cols = len(data), len(data[0])
    gt = s.shapes.add_table(rows, cols, Inches(0.75), Inches(2.05),
                            Inches(11.8), Inches(3.5)).table
    gt.columns[0].width = Inches(4.3)
    gt.columns[1].width = Inches(3.75)
    gt.columns[2].width = Inches(3.75)
    # estilo: quitar banding nativo
    tblPr = gt._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.18); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HD
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TBL_R1 if ri % 2 else TBL_R2
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.name = FONT
            if ri == 0:
                r.font.bold = True; r.font.color.rgb = ACCENT; r.font.size = Pt(16)
            elif ci == 0:
                r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(15)
            elif ci == 1:
                r.font.bold = True; r.font.color.rgb = ACCENT; r.font.size = Pt(15)
            else:
                r.font.color.rgb = MUTE; r.font.size = Pt(15)
    # nota
    tb = s.shapes.add_textbox(Inches(0.75), Inches(5.85), Inches(11.8), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "No es otro modelo: es integracion local y soberana que actua."
    r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_B
    footer(s, 7, total)


# =================================================================
# SLIDE 8 — Arquitectura
# =================================================================
def slide8(total):
    s = add_slide()
    title(s, "Arquitectura", "Modular y escalable")
    bullets(s, [
        [("Núcleo de orquestacion ", True), ("→ despacha a cada capacidad", False)],
        [("Motor LLM local (Ollama) ", True), ("+ router multi-proveedor con failover", False)],
        [("Capacidades aisladas ", True), ("(core/*.py): voz, imagen, sistema, memoria…", False)],
        [("Memoria con recuperacion semantica (RAG)", False)],
    ], top=Inches(2.5), size=19, gap=16)
    # franja de principios
    pn = panel(s, Inches(0.75), Inches(5.7), Inches(11.8), Inches(0.85),
               color=BG2, line=ACCENT_D, line_w=Pt(1))
    tf = pn.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = ("Principios:  local-first  ·  aislamiento  ·  "
              "degradacion gracil  ·  seguridad en la auto-mejora")
    r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = ACCENT; r.font.name = FONT
    footer(s, 8, total)


# =================================================================
# SLIDE 9 — Estado del proyecto (metricas)
# =================================================================
def slide9(total):
    s = add_slide()
    title(s, "Estado del proyecto", "Funcional y en uso real")
    cards = [
        ("~146", "modulos"),
        ("~96.000", "lineas de Python"),
        ("v1.0 → v6.1", "evolucion documentada"),
        ("2", "modelos locales (general + codigo)"),
    ]
    x = Inches(0.75); y = Inches(2.45)
    cw = Inches(2.78); ch = Inches(2.2); gap = Inches(0.23)
    for i, (big, small) in enumerate(cards):
        pn = panel(s, x + i * (cw + gap), y, cw, ch, color=BG2, line=ACCENT)
        tf = pn.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = big
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = FONT_B
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
        r2 = p2.add_run(); r2.text = small
        r2.font.size = Pt(13); r2.font.color.rgb = MUTE; r2.font.name = FONT
    # checks
    bullets(s, [
        [("✅ Suite de pruebas amplia", False)],
    ], top=Inches(5.15), size=18, gap=8)
    footer(s, 9, total)


# =================================================================
# SLIDE 10 — Roadmap
# =================================================================
def slide10(total):
    s = add_slide()
    title(s, "Roadmap: v7.0 “Unbound”",
          "La proxima evolucion, en 4 lineas")
    items = [
        ("1", "\U0001F39B️ Conductor", "orquestar recursos (GPU) + plataforma de tools"),
        ("2", "\U0001F441️ Sentidos en tiempo real", "vision continua"),
        ("3", "\U0001F916 Autonomia", "agente de tareas multi-paso"),
        ("4", "\U0001F4F1 Alcance", "control remoto / movil"),
    ]
    y = Inches(2.4); rh = Inches(0.85); gap = Inches(0.18)
    for i, (num, head, desc) in enumerate(items):
        yy = y + i * (rh + gap)
        # numero circular
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), yy, rh, rh)
        circ.fill.solid(); circ.fill.fore_color.rgb = BG2
        circ.line.color.rgb = ACCENT; circ.line.width = Pt(1.5); _no_shadow(circ)
        ctf = circ.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num
        cr.font.size = Pt(24); cr.font.bold = True; cr.font.color.rgb = ACCENT; cr.font.name = FONT_B
        # texto
        tb = s.shapes.add_textbox(Inches(1.95), yy, Inches(10.7), rh)
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        rh1 = p.add_run(); rh1.text = head + "  —  "
        rh1.font.size = Pt(19); rh1.font.bold = True; rh1.font.color.rgb = WHITE; rh1.font.name = FONT_B
        rh2 = p.add_run(); rh2.text = desc
        rh2.font.size = Pt(17); rh2.font.color.rgb = MUTE; rh2.font.name = FONT
    tb = s.shapes.add_textbox(Inches(0.85), Inches(6.5), Inches(11.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Documento vivo, pensado para escalar."
    r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = ACCENT; r.font.name = FONT
    footer(s, 10, total)


# =================================================================
# SLIDE 11 — Impacto
# =================================================================
def slide11(total):
    s = add_slide()
    title(s, "Impacto", "Por que importa")
    bullets(s, [
        [("\U0001F6E1️ Privacidad y ", True), ("soberania tecnologica", True)],
        [("♿ Accesibilidad ", True), ("(control por voz del equipo)", False)],
        [("\U0001F393 Educacion: ", True), ("IA aplicada y arquitectura real", False)],
        [("⚡ Productividad: ", True), ("automatizacion del escritorio", False)],
        [("\U0001F30E Independencia: ", True), ("sin proveedores externos", False)],
    ], top=Inches(2.45), size=19, gap=14)
    footer(s, 11, total)


# =================================================================
# SLIDE 12 — Cierre
# =================================================================
def slide12(total):
    s = add_slide()
    tech_corners(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.33), Inches(1.5))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "GENESIS"
    r.font.size = Pt(64); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_B
    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.33), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "IA personal, local y soberana."
    r2.font.size = Pt(24); r2.font.color.rgb = ACCENT; r2.font.name = FONT
    tb3 = s.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.33), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "Piensa. Ve. Habla. Crea. Actua. En tu PC."
    r3.font.size = Pt(18); r3.font.bold = True; r3.font.color.rgb = WHITE; r3.font.name = FONT_B
    tb4 = s.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.5))
    p4 = tb4.text_frame.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run(); r4.text = "Alex Quiñones — (contacto)"
    r4.font.size = Pt(14); r4.font.color.rgb = MUTE; r4.font.name = FONT


def main():
    TOTAL = 12
    slide1(TOTAL); slide2(TOTAL); slide3(TOTAL); slide4(TOTAL)
    slide5(TOTAL); slide6(TOTAL); slide7(TOTAL); slide8(TOTAL)
    slide9(TOTAL); slide10(TOTAL); slide11(TOTAL); slide12(TOTAL)
    prs.save(OUT)
    print("PPTX generado:", OUT)
    print("Diapositivas:", len(prs.slides._sldIdLst))
    print("Tamano (bytes):", os.path.getsize(OUT))


if __name__ == "__main__":
    main()
