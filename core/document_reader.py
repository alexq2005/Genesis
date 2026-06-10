"""
GENESIS — Document Reader.

Lee documentos de múltiples formatos y retorna texto plano.
Soporta: PDF, DOCX, PPTX, XLSX, CSV, RTF, TXT, Imágenes (OCR), Audio, Video.
"""
import json
import os
import re
from pathlib import Path
from typing import Optional

# Configurar Tesseract para Windows
_TESSERACT_PATHS = [
    r"C:\Program Files\Tesseract-OCR	esseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR	esseract.exe",
    os.path.expanduser(r"~\AppData\Local\Tesseract-OCR	esseract.exe"),
]
for _tp in _TESSERACT_PATHS:
    if os.path.exists(_tp):
        try:
            import pytesseract as _pt
            _pt.pytesseract.tesseract_cmd = _tp
        except ImportError:
            pass
        break


class DocumentReader:
    """Lee documentos de multiples formatos y retorna texto plano."""

    SUPPORTED_FORMATS = {
        # Documentos
        ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".rtf",
        ".xlsx", ".xls", ".csv",
        ".txt", ".md", ".log", ".json", ".xml", ".yaml", ".yml",
        # Codigo fuente
        ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".css",
        ".java", ".cpp", ".c", ".h", ".cs", ".go", ".rs", ".rb",
        ".php", ".swift", ".kt", ".r", ".sql", ".sh", ".bat",
        # Imagenes (OCR / analisis)
        ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp",
        ".gif", ".ico", ".svg", ".heic", ".heif",
        # Audio (transcripcion)
        ".mp3", ".wav", ".ogg", ".flac", ".aac", ".wma", ".m4a",
        # Video (keyframes + audio)
        ".mp4", ".avi", ".mkv", ".mov", ".webm", ".flv", ".wmv",
    }

    def read(self, filepath: str) -> dict:
        """
        Lee un documento y retorna sus datos.

        Returns:
            dict con keys: text, pages, metadata, tables, format, error
        """
        filepath = str(filepath).strip().strip('"').strip("'")

        if not os.path.exists(filepath):
            return {"error": f"Archivo no encontrado: {filepath}"}

        ext = Path(filepath).suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            return {"error": f"Formato no soportado: {ext}. Soportados: {', '.join(sorted(self.SUPPORTED_FORMATS))}"}

        try:
            if ext == ".pdf":
                return self._read_pdf(filepath)
            elif ext in (".docx", ".doc"):
                return self._read_docx(filepath)
            elif ext in (".pptx", ".ppt"):
                return self._read_pptx(filepath)
            elif ext == ".rtf":
                return self._read_rtf(filepath)
            elif ext in (".xlsx", ".xls"):
                return self._read_xlsx(filepath)
            elif ext == ".csv":
                return self._read_csv(filepath)
            elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp",
                         ".gif", ".ico", ".heic", ".heif"):
                return self._read_image_ocr(filepath)
            elif ext == ".svg":
                return self._read_svg(filepath)
            elif ext in (".mp3", ".wav", ".ogg", ".flac", ".aac", ".wma", ".m4a"):
                return self._read_audio(filepath)
            elif ext in (".mp4", ".avi", ".mkv", ".mov", ".webm", ".flv", ".wmv"):
                return self._read_video(filepath)
            else:
                return self._read_text(filepath)
        except Exception as e:
            return {"error": f"Error leyendo {ext}: {str(e)}"}

    def _read_pdf(self, filepath: str) -> dict:
        """Lee PDF usando PyMuPDF (fitz). Robusto contra PDFs dañados/protegidos."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {"error": "PyMuPDF no instalado. Ejecuta: pip install PyMuPDF"}

        doc = None
        try:
            doc = fitz.open(filepath)

            # Verificar si esta encriptado/protegido
            if doc.is_encrypted:
                try:
                    doc.authenticate("")  # Intentar sin contraseña
                except (RuntimeError, ValueError) as e:
                    doc.close()
                    return {"error": f"PDF protegido con contraseña. No se puede leer: {e}"}

            pages_text = []
            tables = []
            metadata = {}

            # Extraer metadata
            try:
                meta = doc.metadata
                if meta:
                    metadata = {
                        "title": meta.get("title", ""),
                        "author": meta.get("author", ""),
                        "subject": meta.get("subject", ""),
                        "creator": meta.get("creator", ""),
                        "creation_date": meta.get("creationDate", ""),
                        "page_count": doc.page_count,
                    }
            except (RuntimeError, ValueError, KeyError) as e:
                # PDF metadata extraction can fail on malformed PDFs
                metadata = {"page_count": doc.page_count}

            for page_num in range(doc.page_count):
                try:
                    page = doc[page_num]
                    text = page.get_text("text")
                    pages_text.append(text)

                    # Intentar extraer tablas si PyMuPDF >= 1.23
                    try:
                        page_tables = page.find_tables()
                        for t in page_tables:
                            table_data = t.extract()
                            if table_data and len(table_data) > 1:
                                tables.append({
                                    "page": page_num + 1,
                                    "rows": len(table_data),
                                    "cols": len(table_data[0]) if table_data else 0,
                                    "headers": table_data[0] if table_data else [],
                                    "data": table_data[1:] if len(table_data) > 1 else [],
                                })
                    except (AttributeError, RuntimeError, ValueError):
                        pass  # Table extraction not available in this PyMuPDF version
                except Exception as page_err:
                    pages_text.append(f"[Error pagina {page_num + 1}: {page_err}]")

            page_count = doc.page_count  # Guardar ANTES de cerrar
            doc.close()
            doc = None  # Marcar como cerrado
            full_text = "\n\n".join(pages_text)

            # Detectar PDF escaneado (poco texto por pagina) — solo si < 5 paginas (OCR es lento)
            if page_count > 0 and page_count <= 5 and len(full_text.strip()) < 50 * page_count:
                # Intentar OCR como fallback
                ocr_result = self._ocr_pdf_pages(filepath)
                if ocr_result and len(ocr_result) > len(full_text):
                    full_text = ocr_result
                    metadata["ocr_applied"] = True

            return {
                "text": full_text,
                "pages": page_count,
                "metadata": metadata,
                "tables": tables,
                "format": "pdf",
            }

        except Exception as e:
            error_msg = str(e)
            if "document closed" in error_msg.lower():
                return {"error": f"PDF dañado o no se pudo abrir correctamente. Intenta con otro archivo."}
            return {"error": f"Error leyendo PDF: {error_msg}"}
        finally:
            # Siempre cerrar el documento si quedó abierto
            if doc is not None:
                try:
                    doc.close()
                except (RuntimeError, OSError):
                    pass  # Document may already be closed or file handle released

    def _ocr_pdf_pages(self, filepath: str) -> str:
        """OCR para PDFs escaneados (imagen por pagina)."""
        try:
            import fitz
            from PIL import Image
            import pytesseract
            import io

            doc = fitz.open(filepath)
            ocr_texts = []

            for page_num in range(min(doc.page_count, 50)):  # Max 50 paginas OCR
                page = doc[page_num]
                # Renderizar pagina como imagen
                mat = fitz.Matrix(2, 2)  # 2x zoom para mejor OCR
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                text = pytesseract.image_to_string(img, lang="spa+eng")
                if text.strip():
                    ocr_texts.append(text)

            doc.close()
            return "\n\n".join(ocr_texts)
        except (ImportError, OSError, RuntimeError) as e:
            # OCR fallback — fitz, PIL, or pytesseract may not be available or may fail
            return ""

    def _read_docx(self, filepath: str) -> dict:
        """Lee DOCX usando python-docx."""
        try:
            from docx import Document
        except ImportError:
            return {"error": "python-docx no instalado. Ejecuta: pip install python-docx"}

        doc = Document(filepath)
        paragraphs = []
        tables = []
        metadata = {}

        # Extraer metadata
        props = doc.core_properties
        if props:
            metadata = {
                "title": props.title or "",
                "author": props.author or "",
                "created": str(props.created) if props.created else "",
                "modified": str(props.modified) if props.modified else "",
                "category": props.category or "",
            }

        # Extraer parrafos
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Marcar headings
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading", "").strip() or "1"
                    paragraphs.append(f"{'#' * int(level)} {text}")
                else:
                    paragraphs.append(text)

        # Extraer tablas
        for i, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                tables.append({
                    "index": i + 1,
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0,
                    "headers": table_data[0] if table_data else [],
                    "data": table_data[1:] if len(table_data) > 1 else [],
                })

        full_text = "\n\n".join(paragraphs)
        # Estimar paginas (aprox 3000 chars por pagina)
        est_pages = max(1, len(full_text) // 3000)

        return {
            "text": full_text,
            "pages": est_pages,
            "metadata": metadata,
            "tables": tables,
            "format": "docx",
        }

    def _read_xlsx(self, filepath: str) -> dict:
        """Lee XLSX usando openpyxl."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            return {"error": "openpyxl no instalado. Ejecuta: pip install openpyxl"}

        wb = load_workbook(filepath, read_only=True, data_only=True)
        sheets_text = []
        tables = []
        metadata = {
            "sheets": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data = []
            sheet_lines = [f"## Hoja: {sheet_name}"]

            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):  # Ignorar filas vacias
                    rows_data.append(cells)
                    sheet_lines.append(" | ".join(cells))

            if rows_data:
                sheets_text.append("\n".join(sheet_lines))
                tables.append({
                    "sheet": sheet_name,
                    "rows": len(rows_data),
                    "cols": len(rows_data[0]) if rows_data else 0,
                    "headers": rows_data[0] if rows_data else [],
                    "data": rows_data[1:] if len(rows_data) > 1 else [],
                })

        wb.close()
        full_text = "\n\n".join(sheets_text)

        return {
            "text": full_text,
            "pages": len(wb.sheetnames),
            "metadata": metadata,
            "tables": tables,
            "format": "xlsx",
        }

    def _read_csv(self, filepath: str) -> dict:
        """Lee CSV."""
        import csv

        # Detectar encoding
        text = self._read_with_fallback_encoding(filepath)
        if text is None:
            return {"error": f"No se pudo leer el archivo CSV (encoding desconocido)"}

        rows = []
        try:
            import io
            reader = csv.reader(io.StringIO(text))
            for row in reader:
                rows.append(row)
        except (csv.Error, ValueError, TypeError) as e:
            return {"error": f"Error parseando CSV: {e}"}

        # Construir texto legible
        lines = []
        for row in rows:
            lines.append(" | ".join(row))

        tables = []
        if rows:
            tables.append({
                "rows": len(rows),
                "cols": len(rows[0]) if rows else 0,
                "headers": rows[0] if rows else [],
                "data": rows[1:] if len(rows) > 1 else [],
            })

        return {
            "text": "\n".join(lines),
            "pages": 1,
            "metadata": {"row_count": len(rows), "col_count": len(rows[0]) if rows else 0},
            "tables": tables,
            "format": "csv",
        }

    def _read_text(self, filepath: str) -> dict:
        """Lee archivos de texto plano (TXT, MD, LOG, JSON, etc)."""
        text = self._read_with_fallback_encoding(filepath)
        if text is None:
            return {"error": "No se pudo leer el archivo (encoding desconocido)"}

        ext = Path(filepath).suffix.lower()
        est_pages = max(1, len(text) // 3000)

        metadata = {
            "encoding": "utf-8",
            "line_count": text.count("\n") + 1,
        }

        # Si es JSON, intentar parsear para metadata
        if ext == ".json":
            try:
                data = json.loads(text)
                if isinstance(data, dict):
                    metadata["json_keys"] = list(data.keys())[:20]
                elif isinstance(data, list):
                    metadata["json_items"] = len(data)
            except json.JSONDecodeError:
                pass

        return {
            "text": text,
            "pages": est_pages,
            "metadata": metadata,
            "tables": [],
            "format": ext.lstrip("."),
        }

    def _read_image_ocr(self, filepath: str) -> dict:
        """Lee imagen con OCR usando pytesseract. Soporta HEIC/HEIF."""
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            return {"error": "Pillow o pytesseract no instalado. Ejecuta: pip install Pillow pytesseract"}

        try:
            ext = Path(filepath).suffix.lower()
            # HEIC/HEIF requiere pillow-heif
            if ext in (".heic", ".heif"):
                try:
                    import pillow_heif
                    pillow_heif.register_heif_opener()
                except ImportError:
                    return {"error": "pillow-heif no instalado. Ejecuta: pip install pillow-heif"}
            img = Image.open(filepath)
            metadata = {
                "size": f"{img.width}x{img.height}",
                "mode": img.mode,
                "format": img.format,
            }

            # OCR con soporte español + inglés
            try:
                text = pytesseract.image_to_string(img, lang="spa+eng")
            except (RuntimeError, OSError) as e:
                # spa+eng language data not available, fall back to eng only
                try:
                    text = pytesseract.image_to_string(img, lang="eng")
                except (RuntimeError, OSError):
                    # eng language data not available, fall back to default
                    text = pytesseract.image_to_string(img)

            if not text.strip():
                return {
                    "text": "[Imagen sin texto detectable por OCR]",
                    "pages": 1,
                    "metadata": metadata,
                    "tables": [],
                    "format": "image",
                    "warning": "No se detecto texto. Puede necesitar Tesseract-OCR instalado.",
                }

            return {
                "text": text,
                "pages": 1,
                "metadata": metadata,
                "tables": [],
                "format": "image_ocr",
            }
        except (OSError, RuntimeError, ValueError) as e:
            return {"error": f"Error procesando imagen: {e}. Asegurate de tener Tesseract-OCR instalado."}

    def _read_pptx(self, filepath: str) -> dict:
        """Lee PowerPoint PPTX usando python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx no instalado. Ejecuta: pip install python-pptx"}

        prs = Presentation(filepath)
        slides_text = []
        tables = []
        metadata = {
            "slide_count": len(prs.slides),
            "slide_width": str(prs.slide_width),
            "slide_height": str(prs.slide_height),
        }

        for i, slide in enumerate(prs.slides):
            slide_lines = [f"## Slide {i + 1}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        tables.append({
                            "slide": i + 1,
                            "rows": len(table_data),
                            "cols": len(table_data[0]) if table_data else 0,
                            "headers": table_data[0] if table_data else [],
                            "data": table_data[1:] if len(table_data) > 1 else [],
                        })
            # Notas del presentador
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"[Notas: {notes}]")
            slides_text.append("\n".join(slide_lines))

        full_text = "\n\n".join(slides_text)
        return {
            "text": full_text,
            "pages": len(prs.slides),
            "metadata": metadata,
            "tables": tables,
            "format": "pptx",
        }

    def _read_rtf(self, filepath: str) -> dict:
        """Lee RTF (Rich Text Format)."""
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            return {"error": "striprtf no instalado. Ejecuta: pip install striprtf"}

        raw = self._read_with_fallback_encoding(filepath)
        if raw is None:
            return {"error": "No se pudo leer el archivo RTF"}

        text = rtf_to_text(raw)
        return {
            "text": text,
            "pages": max(1, len(text) // 3000),
            "metadata": {"original_size": len(raw)},
            "tables": [],
            "format": "rtf",
        }

    def _read_svg(self, filepath: str) -> dict:
        """Lee SVG extrayendo texto de los elementos XML."""
        try:
            import xml.etree.ElementTree as ET
            tree = ET.parse(filepath)
            root = tree.getroot()

            # Namespace SVG
            ns = {"svg": "http://www.w3.org/2000/svg"}
            texts = []

            # Extraer texto de <text>, <tspan>, <title>, <desc>
            for tag in ["text", "tspan", "title", "desc"]:
                for elem in root.iter(f"{{{ns.get('svg', '')}}}{ tag}"):
                    if elem.text and elem.text.strip():
                        texts.append(elem.text.strip())
                # Sin namespace tambien
                for elem in root.iter(tag):
                    if elem.text and elem.text.strip():
                        t = elem.text.strip()
                        if t not in texts:
                            texts.append(t)

            # Metadata del SVG
            metadata = {
                "width": root.get("width", "unknown"),
                "height": root.get("height", "unknown"),
                "viewBox": root.get("viewBox", ""),
            }

            text = "\n".join(texts) if texts else "[SVG sin texto — imagen vectorial]"
            return {
                "text": text,
                "pages": 1,
                "metadata": metadata,
                "tables": [],
                "format": "svg",
            }
        except (SyntaxError, OSError, ValueError) as e:
            # ET.ParseError is a subclass of SyntaxError
            return {"error": f"Error leyendo SVG: {e}"}

    def _read_audio(self, filepath: str) -> dict:
        """Lee audio: extrae metadata y transcribe con Vosk."""
        import subprocess
        metadata = {}

        # 1. Metadata con mutagen
        try:
            import mutagen
            audio_meta = mutagen.File(filepath)
            if audio_meta:
                metadata["duration_s"] = round(audio_meta.info.length, 1) if hasattr(audio_meta.info, "length") else 0
                metadata["sample_rate"] = getattr(audio_meta.info, "sample_rate", None)
                metadata["channels"] = getattr(audio_meta.info, "channels", None)
                metadata["bitrate"] = getattr(audio_meta.info, "bitrate", None)
                # Tags
                if hasattr(audio_meta, "tags") and audio_meta.tags:
                    for key in ["title", "artist", "album", "genre", "date"]:
                        for tag_key in audio_meta.tags.keys():
                            if key in tag_key.lower():
                                val = audio_meta.tags[tag_key]
                                metadata[key] = str(val[0]) if isinstance(val, list) else str(val)
                                break
        except (ImportError, OSError, ValueError, KeyError) as e:
            pass  # mutagen may not be installed or audio file metadata unreadable

        # 2. Convertir a WAV 16kHz mono para Vosk
        duration = metadata.get("duration_s", 0)
        text = ""
        try:
            ffmpeg_bin = self._get_ffmpeg_path()
            if ffmpeg_bin:
                import tempfile
                wav_tmp = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
                wav_tmp.close()
                try:
                    # Limitar a 10 min para no bloquear
                    time_limit = min(duration, 600) if duration else 600
                    cmd = [
                        ffmpeg_bin, "-y", "-i", filepath,
                        "-ar", "16000", "-ac", "1", "-f", "wav",
                        "-t", str(time_limit),
                        wav_tmp.name
                    ]
                    subprocess.run(cmd, capture_output=True, timeout=120)
                    text = self._transcribe_wav(wav_tmp.name)
                finally:
                    try:
                        os.unlink(wav_tmp.name)
                    except OSError:
                        pass  # Temp file may already be removed or locked
        except (subprocess.SubprocessError, OSError, TimeoutError) as e:
            metadata["transcription_error"] = str(e)

        if not text:
            dur_str = f"{int(duration // 60)}:{int(duration % 60):02d}" if duration else "desconocida"
            text = f"[Audio — Duracion: {dur_str}. No se pudo transcribir. Asegurate de tener Vosk instalado.]"
            if metadata.get("title"):
                text = f"Titulo: {metadata['title']}\n" + text

        metadata["format"] = Path(filepath).suffix.lstrip(".")
        return {
            "text": text,
            "pages": 1,
            "metadata": metadata,
            "tables": [],
            "format": "audio",
        }

    def _read_video(self, filepath: str) -> dict:
        """Lee video: extrae keyframes, metadata y transcribe audio."""
        import subprocess
        metadata = {}
        sections = []

        # 1. Metadata con ffprobe
        ffmpeg_bin = self._get_ffmpeg_path()
        ffprobe_bin = ffmpeg_bin.replace("ffmpeg", "ffprobe") if ffmpeg_bin else None

        try:
            if ffprobe_bin and os.path.exists(ffprobe_bin):
                cmd = [ffprobe_bin, "-v", "quiet", "-print_format", "json",
                       "-show_format", "-show_streams", filepath]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
                if result.stdout:
                    probe = json.loads(result.stdout)
                    fmt = probe.get("format", {})
                    metadata["duration_s"] = round(float(fmt.get("duration", 0)), 1)
                    metadata["size_mb"] = round(int(fmt.get("size", 0)) / 1024 / 1024, 1)
                    metadata["format_name"] = fmt.get("format_long_name", "")
                    for stream in probe.get("streams", []):
                        if stream.get("codec_type") == "video":
                            metadata["resolution"] = f"{stream.get('width', '?')}x{stream.get('height', '?')}"
                            metadata["fps"] = stream.get("r_frame_rate", "?")
                            metadata["video_codec"] = stream.get("codec_name", "?")
                        elif stream.get("codec_type") == "audio":
                            metadata["audio_codec"] = stream.get("codec_name", "?")
                            metadata["audio_channels"] = stream.get("channels", "?")
        except (subprocess.SubprocessError, OSError, TimeoutError, json.JSONDecodeError, ValueError, KeyError) as e:
            pass  # ffprobe may not be available or output may be unparseable

        # 2. Extraer keyframes con OpenCV
        try:
            import cv2
            cap = cv2.VideoCapture(filepath)
            if cap.isOpened():
                total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                fps = cap.get(cv2.CAP_PROP_FPS) or 30
                duration = total_frames / fps if fps > 0 else 0
                if not metadata.get("duration_s"):
                    metadata["duration_s"] = round(duration, 1)
                if not metadata.get("resolution"):
                    w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                    h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                    metadata["resolution"] = f"{w}x{h}"

                # Extraer hasta 10 keyframes distribuidos uniformemente
                num_keyframes = min(10, max(3, int(duration / 30)))  # 1 cada ~30s
                frame_indices = [int(i * total_frames / num_keyframes) for i in range(num_keyframes)]

                keyframe_descriptions = []
                for idx in frame_indices:
                    cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
                    ret, frame = cap.read()
                    if ret:
                        timestamp = idx / fps if fps > 0 else 0
                        ts_str = f"{int(timestamp // 60)}:{int(timestamp % 60):02d}"
                        # Analisis basico del frame
                        avg_brightness = frame.mean()
                        h_frame, w_frame = frame.shape[:2]
                        # Detectar si tiene texto (OCR en keyframe)
                        frame_text = ""
                        try:
                            from PIL import Image
                            import pytesseract
                            import numpy as np
                            rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                            pil_img = Image.fromarray(rgb)
                            frame_text = pytesseract.image_to_string(pil_img, lang="spa+eng").strip()
                        except (ImportError, OSError, RuntimeError):
                            pass  # OCR on video keyframe is best-effort
                        desc = f"[{ts_str}]"
                        if frame_text:
                            desc += f" Texto detectado: {frame_text[:200]}"
                        keyframe_descriptions.append(desc)

                cap.release()
                if keyframe_descriptions:
                    sections.append("## Keyframes\n" + "\n".join(keyframe_descriptions))
        except (ImportError, OSError, RuntimeError, ValueError) as e:
            # OpenCV may not be installed or video file may be unreadable
            metadata["keyframe_error"] = str(e)

        # 3. Transcribir audio del video
        audio_text = ""
        try:
            if ffmpeg_bin:
                import tempfile
                wav_tmp = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
                wav_tmp.close()
                try:
                    duration_val = metadata.get("duration_s", 0)
                    time_limit = min(duration_val, 600) if duration_val else 600
                    cmd = [
                        ffmpeg_bin, "-y", "-i", filepath,
                        "-ar", "16000", "-ac", "1", "-f", "wav",
                        "-t", str(time_limit), "-vn",  # -vn = sin video
                        wav_tmp.name
                    ]
                    subprocess.run(cmd, capture_output=True, timeout=180)
                    audio_text = self._transcribe_wav(wav_tmp.name)
                finally:
                    try:
                        os.unlink(wav_tmp.name)
                    except OSError:
                        pass  # Temp file may already be removed or locked
        except (subprocess.SubprocessError, OSError, TimeoutError) as e:
            metadata["audio_transcription_error"] = str(e)

        if audio_text:
            sections.append("## Transcripcion del audio\n" + audio_text)

        # 4. Construir texto final
        dur = metadata.get("duration_s", 0)
        dur_str = f"{int(dur // 60)}:{int(dur % 60):02d}" if dur else "desconocida"
        header = f"Video: {Path(filepath).name}\nDuracion: {dur_str}"
        if metadata.get("resolution"):
            header += f" | Resolucion: {metadata['resolution']}"

        full_text = header + "\n\n" + "\n\n".join(sections) if sections else header
        if not sections:
            full_text += "\n\n[No se pudo extraer contenido del video]"

        return {
            "text": full_text,
            "pages": 1,
            "metadata": metadata,
            "tables": [],
            "format": "video",
        }

    def _get_ffmpeg_path(self) -> Optional[str]:
        """Obtiene ruta a ffmpeg (imageio_ffmpeg bundle o sistema)."""
        import shutil
        # 1. Sistema
        sys_ffmpeg = shutil.which("ffmpeg")
        if sys_ffmpeg:
            return sys_ffmpeg
        # 2. imageio_ffmpeg bundle
        try:
            import imageio_ffmpeg
            return imageio_ffmpeg.get_ffmpeg_exe()
        except (ImportError, OSError, RuntimeError):
            pass  # imageio_ffmpeg not installed or ffmpeg binary not found
        return None

    def _transcribe_wav(self, wav_path: str) -> str:
        """Transcribe WAV 16kHz mono usando Vosk."""
        try:
            import vosk
            import wave

            # Buscar modelo Vosk
            model_paths = [
                os.path.join(os.path.dirname(__file__), "..", "models", "vosk-model-small-es"),
                os.path.join(os.path.dirname(__file__), "..", "models", "vosk-model-small-es-0.42"),
                os.path.join(os.path.dirname(__file__), "..", "models", "vosk-model-es"),
                os.path.join(os.path.dirname(__file__), "..", "models", "vosk-model-small-en-us-0.15"),
            ]
            # Buscar cualquier carpeta vosk-model en models/
            models_dir = os.path.join(os.path.dirname(__file__), "..", "models")
            if os.path.isdir(models_dir):
                for d in os.listdir(models_dir):
                    full = os.path.join(models_dir, d)
                    if os.path.isdir(full) and "vosk" in d.lower():
                        if full not in model_paths:
                            model_paths.append(full)

            model_path = None
            for mp in model_paths:
                if os.path.isdir(mp):
                    model_path = mp
                    break

            if not model_path:
                return ""

            vosk.SetLogLevel(-1)
            model = vosk.Model(model_path)

            wf = wave.open(wav_path, "rb")
            rec = vosk.KaldiRecognizer(model, wf.getframerate())
            rec.SetWords(True)

            results = []
            while True:
                data = wf.readframes(4000)
                if len(data) == 0:
                    break
                if rec.AcceptWaveform(data):
                    part = json.loads(rec.Result())
                    text = part.get("text", "")
                    if text:
                        results.append(text)

            final = json.loads(rec.FinalResult())
            final_text = final.get("text", "")
            if final_text:
                results.append(final_text)

            wf.close()
            return " ".join(results)
        except (ImportError, OSError, RuntimeError, ValueError) as e:
            # Vosk may not be installed, model not found, or WAV file invalid
            return ""

    def _read_with_fallback_encoding(self, filepath: str) -> Optional[str]:
        """Intenta leer con multiples encodings."""
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]
        for enc in encodings:
            try:
                with open(filepath, "r", encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        return None
